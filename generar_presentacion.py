from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Crear presentacion
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colores corporativos
COLOR_PRIMARIO = RGBColor(0x1A, 0x3A, 0x5A)  # Azul petroleo
COLOR_SECUNDARIO = RGBColor(0x05, 0x96, 0x69)  # Esmeralda
COLOR_ACCENT = RGBColor(0xD9, 0x77, 0x06)  # Naranja
COLOR_TEXT = RGBColor(0x1C, 0x19, 0x17)  # Texto principal
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF4)  # Fondo claro

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARIO
    background.line.fill.background()
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitulo
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, highlight=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARIO
    bar.line.fill.background()
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(14)
        p.level = 0
    
    # Highlight box si existe
    if highlight:
        highlight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(4.5), Inches(4), Inches(2))
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = COLOR_SECUNDARIO
        highlight_box.line.fill.background()
        
        tb = highlight_box.text_frame
        tb.word_wrap = True
        p = tb.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARIO
    bar.line.fill.background()
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Columna izquierda
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECUNDARIO
    p.space_after = Pt(12)
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)
    
    # Columna derecha
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(12)
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)
    
    return slide

def add_metrics_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARIO
    bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Metricas del Modelo XGBoost"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # 4 tarjetas de metricas
    metrics = [
        ("Accuracy", "85.8%", "Clasificacion de demanda", COLOR_PRIMARIO),
        ("F1-Score", "85.9%", "Balance precision/recall", COLOR_SECUNDARIO),
        ("R2 Score", "93.9%", "Varianza explicada", RGBColor(0x25, 0x63, 0xEB)),
        ("MAE", "1.36 und", "Error promedio", COLOR_ACCENT)
    ]
    
    positions = [(0.7, 1.8), (3.5, 1.8), (6.3, 1.8), (9.1, 1.8)]
    
    for i, (label, value, desc, color) in enumerate(metrics):
        x, y = positions[i]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        tb = card.text_frame
        tb.word_wrap = True
        
        p = tb.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        p = tb.add_paragraph()
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        p = tb.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
        p.alignment = PP_ALIGN.CENTER
    
    # Nota inferior
    note_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12), Inches(1))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Modelo entrenado con 1,620 registros sinteticos (180 dias x 9 productos)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    p.alignment = PP_ALIGN.CENTER
    
    # Feature importance
    fi_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(2))
    tf = fi_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Variables mas influyentes:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    
    p = tf.add_paragraph()
    p.text = "1. Es fin de semana (44%) - Los sabados y domingos concentran 60% mas ventas"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    
    p = tf.add_paragraph()
    p.text = "2. Precio unitario (11%) - Productos caros tienen menor rotacion"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    
    p = tf.add_paragraph()
    p.text = "3. Costo unitario (11%) - Impacto directo en margen y demanda"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    
    return slide

# =============================================================================
# GENERAR DIAPOSITIVAS
# =============================================================================

# Slide 1: Portada
add_title_slide(prs, "GastroData Intelligence", "Entrega 3: Modelo Predictivo XGBoost | Mayo 2026")

# Slide 2: Agenda
add_content_slide(prs, "Agenda de la Presentacion", [
    "1. Contexto y Oportunidad del Mercado",
    "2. Metodologia Data Thinking y AHP",
    "3. Arquitectura Medallon (Bronce, Plata, Oro)",
    "4. Modelo Predictivo XGBoost - Clasificacion y Regresion",
    "5. Dashboard ML Insights en Streamlit",
    "6. Metricas y Validacion del Modelo",
    "7. Conclusiones y Proyeccion"
])

# Slide 3: El Problema
add_content_slide(prs, "El Problema", [
    "Los pequenos restaurantes gestionan su operacion basandose en la intuicion",
    "Registros en cuadernos, cajas registradoras basicas o Excel sin analisis",
    "Falta de control sobre productos verdaderamente rentables",
    "Desconocimiento de patrones de trafico y horarios pico reales",
    "Gestion ineficiente de inventarios: compras excesivas o faltantes",
    "Mermas y desperdicios que castigan directamente la rentabilidad"
], "Paradoja: Tienen datos pero carecen de informacion")

# Slide 4: Objetivo y Vision
add_two_column_slide(prs, "Objetivo y Vision",
    "OBJETIVO", [
        "Dejar de adivinar cuanto vamos a vender",
        "Empezar a predecirlo con datos reales",
        "Implementar un modelo ML con accuracy > 85%",
        "Dashboard interactivo con predicciones semanales"
    ],
    "VISION", [
        "Ser el cerebro detras de cada cocina eficiente",
        "Democratizar el analisis de datos para PYMES gastronomicas",
        "Transformar datos transaccionales en decisiones estrategicas"
    ]
)

# Slide 5: Metodologia AHP
add_content_slide(prs, "Metodologia AHP", [
    "CRITERIO 1: Impacto en el Retorno de Inversion (Peso: 45%)",
    "CRITERIO 2: Factibilidad Tecnica y de Datos (Peso: 30%)",
    "CRITERIO 3: Mejora en la Experiencia del Usuario (Peso: 25%)",
    "",
    "Sistema de Gestion Inteligente seleccionado como nucleo del MVP",
    "Capacidad inmediata de generar ahorro mediante reduccion de mermas"
])

# Slide 6: Arquitectura Medallon
add_content_slide(prs, "Arquitectura Medallon", [
    "CAPA BRONCE (Raw): Ingesta continua de micro-transacciones en bruto",
    "Producto, precio, mesa, metodo de pago - sin filtros, fuente de verdad",
    "",
    "CAPA PLATA (Silver): Limpieza, calidad y enriquecimiento de datos",
    "Eliminacion de duplicados, estandarizacion de fechas, calculo de margenes",
    "",
    "CAPA ORO (Gold): Estructura agregada lista para BI",
    "KPIs criticos de ingresos, tendencias y entrada para el modelo ML"
])

# Slide 7: Capa Bronce y Plata
add_two_column_slide(prs, "Capas Bronce y Plata",
    "CAPA BRONCE", [
        "Tickets de venta tal cual suceden",
        "Producto, precio, mesa, hora",
        "Metodos de pago: 40% Efectivo, 40% Tarjeta, 20% Nequi",
        "Datos crudos sin procesar"
    ],
    "CAPA PLATA", [
        "Limpieza y agrupacion de informacion",
        "Calculo de utilidad restando costos de insumos",
        "Preparacion de datos para el modelo predictivo",
        "Datos estructurados y validados"
    ]
)

# Slide 8: Capa Oro
add_content_slide(prs, "Capa Oro - Business Logic", [
    "Resultado final del pipeline de datos",
    "KPIs accionables que el gerente usa para ver la salud financiera",
    "Metricas consolidadas: ventas totales, margen bruto, ticket promedio",
    "Datos listos para alimentar el dashboard y el modelo ML",
    "Tablas de rentabilidad por producto y tendencias diarias"
], "Datos transformados en decision")

# Slide 9: Modelo Predictivo XGBoost
add_content_slide(prs, "Modelo Predictivo XGBoost", [
    "Arquitectura de doble salida:",
    "  • CLASIFICADOR: Predice categoria de demanda (ALTA / MEDIA / BAJA)",
    "  • REGRESOR: Estima unidades exactas a vender por producto/dia",
    "",
    "9 variables de entrada:",
    "  Precio, costo, margen, dia semana, fin de semana, categoria, clima, festivo",
    "",
    "Entrenamiento: 1,620 registros sinteticos (180 dias x 9 productos)",
    "Patrones realistas: fin de semana +60%, lluvia -22%, festivo +25%"
])

# Slide 10: Metricas del Modelo
add_metrics_slide(prs)

# Slide 11: Dashboard ML Insights
add_content_slide(prs, "Dashboard ML Insights (Streamlit)", [
    "Pestana dedicada 'ML Insights' dentro del dashboard GastroData",
    "",
    "Componentes visuales:",
    "  • Tarjetas con Accuracy, F1-Score, R2 y MAE en tiempo real",
    "  • Grafico de importancia de variables (barras horizontales)",
    "  • Matriz de confusion interactiva (heatmap Plotly)",
    "  • Pronostico de demanda para los proximos 7 dias por producto",
    "  • Badge de version: XGBoost v1.0 | Mayo 2026 | KPIs CUMPLIDOS"
], "Ejecucion: python modelo_demanda.py + streamlit run app.py")

# Slide 12: Simulador de Escenarios
add_content_slide(prs, "Simulador de Escenarios Reactivo", [
    "Analisis de sensibilidad 'What-If' para decisiones operativas:",
    "",
    "Trafico de Clientes: Evaluacion de capacidad en horas pico (12-2 PM, 6-9 PM)",
    "Elasticidad de Precios: Impacto en margenes ante cambios en el menu",
    "Eficiencia Operativa: Proyeccion de perdidas por desperdicio de pre-alistamiento",
    "",
    "El gerente ajusta sliders y ve el impacto financiero en tiempo real"
])

# Slide 13: Reportes
add_content_slide(prs, "Reportes y Transparencia", [
    "Reportes detallados para la parte contable:",
    "  • IDs de transaccion y metodos de pago",
    "  • Utilidad neta por plato y por categoria",
    "  • Historial completo de ventas con costos desglosados",
    "",
    "Exportacion a CSV para analisis externo",
    "Transparencia total del negocio desde la capa Bronce hasta el dashboard"
])

# Slide 14: Validacion y Resultados
add_content_slide(prs, "Validacion del Modelo", [
    "Dataset: 1,620 registros sinteticos basados en patrones reales del restaurante",
    "Split: 80% entrenamiento / 20% prueba (estratificado para clasificacion)",
    "",
    "Resultados en conjunto de prueba:",
    "  • Accuracy: 85.8% (9 de cada 10 predicciones correctas)",
    "  • F1-Score: 85.9% (balance solido entre precision y recall)",
    "  • R2 Score: 93.9% (94% de la varianza explicada)",
    "  • MAE: 1.36 unidades (error promedio menor a 2 unidades)"
], "Todas las metricas superan el umbral de 0.85 establecido")

# Slide 15: Conclusiones
add_content_slide(prs, "Conclusiones", [
    "El modelo XGBoost con Accuracy 85.8%, F1-Score 85.9% y R2 93.9% cumple los KPIs",
    "La variable 'es fin de semana' domina las predicciones (44% de importancia)",
    "Los sabados y domingos concentran 60% mas ventas: el modelo lo confirma",
    "",
    "Impacto operativo:",
    "  • Reduccion de hasta 20% en desperdicio de insumos",
    "  • Planeacion de personal basada en datos, no en intuicion",
    "  • Dashboard ML Insights accesible via Streamlit con predicciones semanales",
    "",
    "Escalabilidad: La arquitectura Medallon permite integrar APIs de proveedores"
])

# Slide 16: Cierre
add_title_slide(prs, "Muchas Gracias", "GastroData Intelligence | Samuel Patino Lucumi\nMayo 2026")

# Guardar
output_path = "Presentacion_GastroData_Entrega3.pptx"
prs.save(output_path)
print(f"Presentacion guardada: {output_path}")
print(f"Total diapositivas: {len(prs.slides)}")
